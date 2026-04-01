"""
make_week1_overview.py — Generate the Week 1 PowerPoint overview

Usage:
    pip install python-pptx
    python make_week1_overview.py

Output:
    week1-overview.pptx in the current directory.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_content_slide(prs, title, bullets, notes=""):
    """Add a slide with a title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    # Format all paragraphs
    for p in tf.paragraphs:
        p.font.size = Pt(16)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_two_column_slide(prs, title, left_items, right_items, left_title="", right_title=""):
    """Add a slide with two columns of content."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    # Left column
    left_top = Inches(1.3)
    txBox_l = slide.shapes.add_textbox(Inches(0.5), left_top, Inches(4.3), Inches(5))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    if left_title:
        tf_l.text = left_title
        tf_l.paragraphs[0].font.bold = True
        tf_l.paragraphs[0].font.size = Pt(16)
    for item in left_items:
        p = tf_l.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)

    # Right column
    txBox_r = slide.shapes.add_textbox(Inches(5.2), left_top, Inches(4.3), Inches(5))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    if right_title:
        tf_r.text = right_title
        tf_r.paragraphs[0].font.bold = True
        tf_r.paragraphs[0].font.size = Pt(16)
    for item in right_items:
        p = tf_r.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.3), Inches(9), Inches(0.4 * num_rows)
    )
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    add_title_slide(
        prs,
        "Week 1: Discrete Event Systems\n& Supervisory Control",
        "PhD Preparation — Holonic Demanufacturing\nunder Structural Uncertainty"
    )

    # --- Slide 2: Week Overview ---
    add_content_slide(prs, "Week 1 — Overview", [
        "Day 1: DES foundations — what is a DES, event alphabets",
        "Day 2: Automata for DES — DFA, marked states, languages",
        "Day 3: Supervisory Control Theory — plant/supervisor split",
        "Day 4: Controllability & Observability",
        "Day 5: Safety, Liveness, Blocking/Nonblocking",
        "Day 6: Complete worked example — automaton + supervisor",
        "Day 7: Introduction to Petri Nets",
    ])

    # --- Slide 3: What is a DES? ---
    add_content_slide(prs, "What is a Discrete-Event System?", [
        "State changes at discrete instants, triggered by events",
        "Not time-driven — the system waits for something to happen",
        "Natural fit for manufacturing/demanufacturing cells:",
        "  → arrivals, inspections, faults, routing decisions",
        "Key modelling primitive: the EVENT",
        "",
        "Source: Cassandras, Discrete Event Systems (EOLSS)",
    ])

    # --- Slide 4: Finite Automata ---
    add_content_slide(prs, "Finite Automata for DES", [
        "Plant automaton: G = (Q, Σ, δ, q₀, Qₘ)",
        "Q = finite set of states",
        "Σ = event alphabet",
        "δ = transition function (state × event → state)",
        "q₀ = initial state",
        "Qₘ = marked states (task completion)",
        "Generated language L(G): all possible event strings",
        "Marked language Lₘ(G): strings ending in marked states",
        "",
        "Source: Raisch, DES & Hybrid Systems, §4.5.1",
    ])

    # --- Slide 5: Plant/Supervisor Loop ---
    add_two_column_slide(
        prs, "The Plant / Supervisor Feedback Loop",
        left_items=[
            "Plant G models ALL possible behaviour",
            "Including unsafe transitions",
            "Generates events (some controllable, some not)",
        ],
        right_items=[
            "Supervisor S restricts behaviour",
            "Observes events",
            "Disables controllable events",
            "Cannot prevent uncontrollable events",
        ],
        left_title="Plant G",
        right_title="Supervisor S",
    )

    # --- Slide 6: Controllable vs Uncontrollable ---
    add_table_slide(
        prs, "Controllable vs Uncontrollable Events",
        ["Property", "Controllable (Σ_c)", "Uncontrollable (Σ_uc)"],
        [
            ["Supervisor can...", "Disable / prevent", "NOT prevent"],
            ["Typical source", "Actuator commands", "Sensor readings, faults"],
            ["Example", "unscrew_cover, route_recycle", "arrival, inspect_ok, fault"],
            ["Engineering rule", "Commands you issue", "Outcomes you observe"],
        ]
    )

    # --- Slide 7: Observability ---
    add_content_slide(prs, "Observability & Partial Observation", [
        "Observable events (Σ_o): supervisor sees them happen",
        "Unobservable events: happen silently",
        "Natural projection P: Σ* → Σ_o* (erases unobservable events)",
        "",
        "Challenge: if two strings look the same under P,",
        "  the supervisor must make the SAME control decision",
        "",
        "Practical impact: hidden damage, unconfirmed actuator commands",
        "",
        "Source: Cai & Wonham (2020), pp. 8–9",
    ])

    # --- Slide 8: Safety vs Liveness ---
    add_two_column_slide(
        prs, "Safety vs Liveness",
        left_items=[
            '"Nothing bad ever happens"',
            "System never enters forbidden states",
            "Enforced by disabling unsafe transitions",
            "Example: never recycle without battery removal",
        ],
        right_items=[
            '"Something good eventually happens"',
            "System can always make progress",
            "Can always reach a marked (completion) state",
            "Example: unit can always be processed to done",
        ],
        left_title="Safety",
        right_title="Liveness / Nonblocking",
    )

    # --- Slide 9: Blocking ---
    add_content_slide(prs, "Blocking, Nonblocking, Deadlock, Livelock", [
        "Nonblocking: every reachable state can reach a marked state",
        "  Formally: L̄ₘ(G) = L(G)",
        "",
        "Deadlock: no transitions enabled (system is stuck)",
        "Livelock: transitions exist but no path to marked state",
        "",
        "Reachable: forward from q₀",
        "Coreachable: backward from Qₘ",
        "Trim = reachable ∩ coreachable",
        "",
        "Source: Raisch, p. 75; Cai & Wonham (2020), p. 2",
    ])

    # --- Slide 10: Worked Example — Plant ---
    add_content_slide(prs, "Worked Example: Demanufacturing Cell — Plant", [
        "Scenario: laptop intake → inspection → open/battery → recycle/quarantine",
        "",
        "States: Idle, Intake, Inspected_OK, Inspected_Suspect,",
        "        Opened, Battery_Removed, Recycle_Done, Quarantine_Done",
        "",
        "Σ_c = {unscrew_cover, remove_battery, route_recycle, route_quarantine}",
        "Σ_uc = {arrival, inspect_ok, inspect_sus}",
        "",
        "Marked states = {Recycle_Done, Quarantine_Done}",
    ])

    # --- Slide 11: Worked Example — Supervisor ---
    add_table_slide(
        prs, "Worked Example: Supervisor Rule Table",
        ["State", "Enabled Σ_c", "Disabled Σ_c", "Rule"],
        [
            ["Inspected_OK", "unscrew_cover, route_quarantine", "route_recycle", "S2"],
            ["Inspected_Suspect", "route_quarantine", "unscrew_cover, route_recycle", "S1"],
            ["Opened", "remove_battery, route_quarantine", "route_recycle", "S2"],
            ["Battery_Removed", "route_recycle, route_quarantine", "—", "—"],
        ]
    )

    # --- Slide 12: Petri Nets Introduction ---
    add_content_slide(prs, "Introduction to Petri Nets", [
        "Bipartite graph: places (circles) + transitions (rectangles)",
        "Tokens in places represent state / resource availability",
        "Firing: transition consumes input tokens, produces output tokens",
        "",
        "Advantages over automata:",
        "  → Natural resource contention modelling (robot, station)",
        "  → Concurrency without state explosion",
        "  → Place invariants for structural safety",
        "",
        "Source: Murata, Petri Nets (1989 tutorial paper)",
    ])

    # --- Slide 13: Petri Net Example ---
    add_content_slide(prs, "Petri Net: Demanufacturing Cell", [
        "Flow places: Unit_In → Wait_Inspect → OK/SUS → OPEN → ...",
        "Resource places: INSPECTOR_FREE, ROBOT_FREE",
        "Status places: SAFE_FLAG/HAZARD_FLAG, BATT_IN/BATT_OUT",
        "",
        "Key invariants:",
        "  INSPECTOR_FREE = 1  (always available)",
        "  ROBOT_FREE = 1  (always available)",
        "  SAFE_FLAG + HAZARD_FLAG = 1  (boolean)",
        "  BATT_IN + BATT_OUT = 1  (boolean)",
        "",
        "Supervisor disables: unscrew_hazard, recycle_direct",
    ])

    # --- Slide 14: Key Takeaways ---
    add_content_slide(prs, "Week 1 — Key Takeaways", [
        "1. Demanufacturing cells are naturally DES (event-driven)",
        "2. Plant automaton models ALL physically possible behaviour",
        "3. Supervisor disables controllable events to enforce safety",
        "4. Uncontrollable events constrain what any supervisor can achieve",
        "5. Nonblocking guarantees the system can always complete safely",
        "6. Petri nets complement automata for resource/concurrency modelling",
        "7. The supervisor-enabled set is the formal foundation",
        "   for all higher layers (twin, learning, mediation)",
    ])

    # --- Slide 15: Next Steps ---
    add_content_slide(prs, "Recommended Next Steps — Week 2", [
        "Week 2 focus: Structural Uncertainty & Belief States",
        "",
        "→ What happens when events are partially observable?",
        "→ Belief states: distributions over possible states",
        "→ Information-state supervisors",
        "→ Connection to the digital twin's state estimation",
        "",
        "Tools: same cell model, extended with unobservable events",
    ])

    return prs


def main():
    prs = build_presentation()
    output_path = "week1-overview.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
